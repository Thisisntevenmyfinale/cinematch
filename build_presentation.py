"""Build the final CineMatch presentation with real metrics and figures.

Covers ALL assignment requirements:
- Technical challenges (explicit)
- Method comparison (direct, with numbers)
- Final remarks (dedicated closing)
- Accuracy is not enough (diversity, novelty, bias, scalability)
- Ethics & social implications
- Development methodology
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RESULTS = Path("results")
FIGURES = RESULTS / "figures"
OUT = Path("CineMatch_Presentation.pptx")

# Colours — matching the CineMatch Streamlit UI
BG      = RGBColor(0x14, 0x14, 0x14)
SURFACE = RGBColor(0x1C, 0x1C, 0x1C)
ACCENT  = RGBColor(0xE8, 0x40, 0x3E)
WHITE   = RGBColor(0xE5, 0xE5, 0xE5)
GREY    = RGBColor(0x80, 0x80, 0x80)
DIM     = RGBColor(0x4A, 0x4A, 0x4A)

W = Inches(13.333)
H = Inches(7.5)


def set_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, title, bullets, start_top=Inches(1.8)):
    set_bg(slide)
    add_text(slide, title, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
             font_size=36, bold=True, color=WHITE)
    # Accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    y = start_top
    for bullet in bullets:
        add_text(slide, bullet, Inches(1.0), y, Inches(10.5), Inches(0.5),
                 font_size=16, color=WHITE)
        y += Inches(0.45)


def add_image_slide(slide, title, img_path, caption=""):
    set_bg(slide)
    add_text(slide, title, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             font_size=32, bold=True, color=WHITE)
    if os.path.exists(img_path):
        slide.shapes.add_picture(str(img_path), Inches(0.8), Inches(1.5),
                                  Inches(11.5), Inches(5.3))
    if caption:
        add_text(slide, caption, Inches(0.8), Inches(6.9), Inches(11), Inches(0.5),
                 font_size=12, color=GREY)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]  # blank layout

    # ── Slide 1: Title ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_text(s, "CINEMATCH", Inches(1), Inches(2), Inches(11), Inches(1.5),
             font_size=72, bold=True, color=ACCENT)
    add_text(s, "Movie Recommender System", Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             font_size=28, color=WHITE)
    add_text(s, "Recommender Systems Course  |  Prof. Marc Torrens  |  ESADE 2025",
             Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             font_size=16, color=GREY)
    add_text(s, "Jan Philipp Gnau", Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             font_size=18, color=WHITE)

    # ── Slide 2: Agenda ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "AGENDA", [
        "1.  Problem & Dataset",
        "2.  Non-Personalized Baselines",
        "3.  Content-Based Filtering (TF-IDF)",
        "4.  Collaborative Filtering (User-User & Item-Item)",
        "5.  Matrix Factorization (Biased SGD)",
        "6.  Evaluation: Accuracy, Beyond Accuracy, Trade-offs",
        "7.  Technical Challenges & Limitations",
        "8.  Ethical & Social Implications",
        "9.  UX & Prototype Demo",
        "10. Conclusions, Future Work & Final Remarks",
    ])

    # ── Slide 3: Problem & Dataset ──────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "PROBLEM & DATASET", [
        "Goal: recommend movies a user will enjoy, from a catalog of 9,742 titles",
        "Dataset: MovieLens Latest Small -- GroupLens Research (University of Minnesota)",
        "  License: research & educational use permitted (see dataset README)",
        "100,836 ratings  |  610 users  |  9,742 movies  |  0.5--5.0 scale",
        "Data representation: User-Item matrix (610 x 9742), 98.3% sparse",
        "No filtering or sampling applied -- all raw data used as-is",
        "Train/Test split: 80/20 random, stratified by user (random_state=42)",
        "",
        "Same dataset, 8 algorithms: enables direct, fair method comparison",
        "Agile approach: basic evaluation running from week 1, extended iteratively",
    ])

    # ── Slide 4: EDA ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_image_slide(s, "EXPLORATORY DATA ANALYSIS",
                    FIGURES / "eda_distributions.png",
                    "Rating distribution, user activity, and item popularity from MovieLens Latest Small.")

    # ── Slide 5: Non-Personalized Baselines ─────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "NON-PERSONALIZED BASELINES", [
        "Most Popular: rank by interaction count (# ratings)",
        "Highest Average: S(u,i) = mean of all ratings for item i",
        "  -- Minimum 20 ratings filter to avoid noise (CollaborativeFiltering.pdf, Folie 12)",
        "Random: control baseline (EvaluationRecommenderSystems.pdf, Folie 28)",
        "",
        "Results: Most Popular achieves Precision@10 = 0.122 (best overall!)",
        "  Strength: crowd favorites are a strong signal in sparse data",
        "  Weakness: 100% popularity bias, 0% serendipity, coverage < 1%",
        "  Use case: cold-start fallback when no user history is available",
    ])

    # ── Slide 6: Content-Based Filtering ────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "CONTENT-BASED FILTERING", [
        "Item representation: TF-IDF on genres (ContentBasedFiltering.pdf, Folie 37)",
        "  -- Optionally enriched with user-generated tags (Folie 17)",
        "User profile: profile(u) = Sigma_i (r_ui - r_u) * vector(i), normalized",
        "  -- (ContentBasedFiltering.pdf, Folien 27-28)",
        "Prediction: score(u,i) = cos(profile_u, vector_i) (Folie 33)",
        "",
        "Results: P@10 = 0.009  |  Coverage = 20.6%  |  Diversity = 0.09",
        "  Strength: best coverage (20.6%) and lowest pop bias (18.6%) -- good for discovery",
        "  Weakness: low accuracy, limited by genre granularity (only 20 genres)",
        "  Use case: niche discovery, cold-start items (no ratings needed for items)",
    ])

    # ── Slide 7: Collaborative Filtering ────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "COLLABORATIVE FILTERING", [
        "User-User CF (CollaborativeFiltering.pdf, Folie 15):",
        "  S(u,i) = r_u + Sigma_v (r_vi - r_v) * w_uv / Sigma |w_uv|",
        "  Similarity: Pearson Correlation (Folie 18), top-k=20 neighbours",
        "",
        "Item-Item CF (CollaborativeFiltering.pdf, Folie 29):",
        "  S(u,i) = r_i + Sigma_j (r_uj - r_j) * w_ij / Sigma |w_ij|",
        "  Similarity: Adjusted Cosine (Folie 26), top-k=20 neighbours",
        "",
        "  Strength: RMSE = 0.893 (User-User), highest novelty = 9.2, 0% pop bias",
        "  Weakness: near-zero P@10 -- 98.3% sparsity means too few co-rated items",
        "  Use case: rating prediction (not top-N) on denser datasets",
    ])

    # ── Slide 8: Matrix Factorization ───────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "MATRIX FACTORIZATION (Biased SGD)", [
        "MatrixFactorization.pdf, Folien 6-7 + Koren et al. (2009):",
        "  r_hat = mu + b_u + b_i + p_u * q_i",
        "  50 latent factors, 20 epochs, lr=0.005, reg=0.02",
        "",
        "SGD: b_u += alpha*(e - lambda*b_u), p_u += alpha*(e*q_i - lambda*p_u)",
        "  Correct p_u_old for q_i update (simultaneous gradient)",
        "",
        "Results: P@10 = 0.053  |  RMSE = 0.881  |  Hit Rate = 32.4%",
        "  Strength: best personalized method (P@10 4x higher than Content-Based)",
        "  Weakness: only 32.4% of users receive at least one relevant hit",
        "  Use case: personalized ranking at scale (O(n_ratings * k) per epoch)",
    ])

    # ── Slide 9: MF Training Curve ──────────────────────────────
    s = prs.slides.add_slide(blank)
    add_image_slide(s, "MF TRAINING CONVERGENCE",
                    FIGURES / "mf_training_curve.png",
                    "Training RMSE decreases steadily over 20 epochs, confirming proper SGD convergence.")

    # ── Slide 10: Accuracy Metrics ──────────────────────────────
    s = prs.slides.add_slide(blank)
    add_image_slide(s, "EVALUATION: ACCURACY METRICS",
                    FIGURES / "model_comparison.png",
                    "Prediction metrics (MAE, RMSE) and ranking metrics (Precision@10, Recall@10, NDCG@10, MRR, Hit Rate@10) across all 8 models.")

    # ── Slide 11: Beyond Accuracy ───────────────────────────────
    s = prs.slides.add_slide(blank)
    add_image_slide(s, "EVALUATION: BEYOND ACCURACY",
                    FIGURES / "per_user_distributions.png",
                    "Beyond-accuracy measures: Coverage, Diversity, Novelty, Popularity Bias, Serendipity, Fairness. "
                    "Per-user distributions reveal high variance -- mean alone does not tell the full story.")

    # ── Slide 12: Method Comparison Table ───────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "METHOD COMPARISON: ALL MODELS, SAME DATA", [
        "                         P@10     RMSE    Coverage  Novelty  Pop Bias   Time",
        "  Most Popular           0.122     --       0.6%      1.7      100%    0.003s",
        "  Highest Average        0.046     --       0.3%      3.4       88%    0.005s",
        "  Content-Based          0.009     --      20.6%      6.8       19%    0.034s",
        "  CB + Tags              0.008     --      21.2%      7.1       16%    0.322s",
        "  User-User CF           0.000    0.893     2.1%      9.1        0%   12.35s",
        "  Item-Item CF           0.000    0.935     5.0%      9.2        0%    3.84s",
        "  Matrix Factorization   0.053    0.881     2.4%      3.3       86%   19.87s",
        "",
        "Ranking metrics | Prediction metrics | Beyond-accuracy measures -- all three categories shown.",
        "No single model dominates all dimensions: this IS the core lecture insight.",
    ])

    # ── Slide 13: Trade-off Analysis ────────────────────────────
    s = prs.slides.add_slide(blank)
    add_image_slide(s, "TRADE-OFF: PRECISION vs DIVERSITY vs NOVELTY",
                    FIGURES / "tradeoff_analysis.png",
                    "Example: Most Popular has 13x higher Precision than Content-Based (0.122 vs 0.009) "
                    "but 34x lower Coverage (0.6% vs 20.6%) and 4x higher Popularity Bias (100% vs 19%). "
                    "As the lectures emphasise: accuracy is not enough.")

    # ── Slide 14: Radar Chart ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_image_slide(s, "MULTI-DIMENSIONAL MODEL COMPARISON",
                    FIGURES / "radar_chart.png",
                    "Normalized radar chart: each model excels in different dimensions. "
                    "Diversity, novelty, bias, and scalability vary as much as accuracy.")

    # ── Slide 15: Fairness & Popularity Bias ────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "FAIRNESS & POPULARITY BIAS", [
        "Fairness: NDCG@K gap between heavy raters (>50th percentile) and light raters",
        "",
        "Largest gap: Most Popular (0.156) -- heavy raters benefit disproportionately",
        "Matrix Factorization gap: 0.143 -- also biased toward active users",
        "CB+Tags gap: 0.005 -- most equitable across user groups",
        "Random gap: 0.003 -- naturally fair (no personalization)",
        "",
        "Popularity bias: Most Popular = 100% from top 10% items",
        "CF methods = 0% popularity bias (but also near-zero precision)",
        "Production systems need explicit debiasing strategies",
    ])

    # ── Slide 16: Technical Challenges & Limitations ────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "TECHNICAL CHALLENGES & LIMITATIONS", [
        "Sparsity (98.3%): the dominant challenge -- User-User CF finds too few",
        "  co-rated items for reliable Pearson correlations, P@10 drops to ~0",
        "",
        "Cold-start: new users without history receive non-personalized fallback;",
        "  new items without ratings are invisible to CF/MF (content-based can help)",
        "",
        "Scalability: User-User CF = O(n_users^2 * n_items), 12.3s to train;",
        "  Most Popular = O(n_ratings), 0.003s -- 4000x faster for higher accuracy",
        "",
        "Popularity bias: Most Popular recommends only the top 0.6% of items,",
        "  creating a rich-get-richer feedback loop in production systems",
    ])

    # ── Slide 17: Ethics & Social Implications ──────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "ETHICAL & SOCIAL IMPLICATIONS", [
        "Filter bubbles: Content-Based has diversity of only 0.09 -- users see",
        "  increasingly narrow genre profiles, limiting cultural exploration",
        "",
        "Popularity bias amplification: Most Popular (100% pop bias) surfaces",
        "  the same 60 movies to all users, marginalizing niche content",
        "",
        "Fairness: heavy raters get 2x better NDCG than light raters (Most Pop),",
        "  disadvantaging casual users who contribute fewer ratings",
        "",
        "Privacy: CF requires storing and processing individual rating histories",
        "Transparency: MF latent factors are not interpretable -- users cannot",
        "  understand why a recommendation was made",
        "User autonomy: users have no control over recommendation logic in this prototype",
    ])

    # ── Slide 18: UX & Demo ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "UX & PROTOTYPE", [
        "Netflix-inspired Streamlit UI with 6 pages",
        "Design: dark theme (#141414), consistent with CineMatch brand identity",
        "Poster lazy loading via TMDb API with JSON cache",
        "",
        "Pages: Home (hero + trending), Discover (search + filter + similar),",
        "  Recommendations (7 algorithms), User Profile (stats + genre radar),",
        "  Compare (side-by-side with overlap heatmap), Evaluation (full metrics)",
        "",
        "Poster cards with hover overlay showing title, year, genres, and 'why'",
        "Horizontal scroll rows mimicking Netflix filmstrip UX",
    ])

    # ── Slide 19: Future Work ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "FUTURE WORK", [
        "Hybrid model (CF + Content-Based): combine MF's personalized ranking",
        "  with CB's item features to address both cold-start and sparsity",
        "",
        "Online learning: update latent factors incrementally as new ratings",
        "  arrive, instead of full retraining (critical for production deployment)",
        "",
        "A/B testing framework: measure real user engagement (clicks, watch time),",
        "  not just offline metrics -- offline P@10 does not predict user satisfaction",
        "",
        "Explicit debiasing: re-rank to enforce minimum diversity/coverage constraints",
        "All four target the identified weaknesses: cold-start, sparsity, pop bias, offline gap",
    ])

    # ── Slide 20: Final Remarks ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bullet_slide(s, "FINAL REMARKS", [
        "All formulas implemented exactly as in Prof. Torrens' lecture slides (verified audit)",
        "Most Popular is surprisingly strong in sparse data (P@10 = 0.122)",
        "  -- but only 32.4% of users get a relevant hit even from best personalized method (MF)",
        "",
        "Core insight: 'accuracy is not enough' (as lectures emphasise)",
        "  -- diversity, novelty, bias, and scalability matter equally in production",
        "  -- no single model dominates; ensembles are the path forward (Netflix Prize, Folie 22)",
        "",
        "When is MF worth 20s training over 0.003s Most Popular? Only when the platform is",
        "  large enough that a 5% precision lift translates to significant engagement gains",
        "",
        "Honest limitation: at 98.3% sparsity, all methods struggle with top-N --",
        "  academically rigorous, but production deployment requires hybrid extensions",
    ])

    prs.save(str(OUT))
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
